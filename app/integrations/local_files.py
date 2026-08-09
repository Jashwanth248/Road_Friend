from __future__ import annotations

import csv
import io
from pathlib import Path


class UnsupportedDocumentError(ValueError):
    pass


class LocalDocumentParser:
    """Parse only files explicitly uploaded through the browser file picker."""

    SUPPORTED = {".pdf", ".txt", ".md", ".csv", ".json", ".docx", ".xlsx", ".pptx"}

    def parse(self, filename: str, data: bytes) -> str:
        suffix = Path(filename).suffix.lower()
        if suffix not in self.SUPPORTED:
            raise UnsupportedDocumentError(
                f"Unsupported file type {suffix or '(none)'}. Supported: {', '.join(sorted(self.SUPPORTED))}"
            )
        if suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)
        if suffix in {".txt", ".md", ".json"}:
            return data.decode("utf-8", errors="ignore")
        if suffix == ".csv":
            text = data.decode("utf-8", errors="ignore")
            rows = csv.reader(io.StringIO(text))
            return "\n".join(" | ".join(row) for row in rows)
        if suffix == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)
        if suffix == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            output = []
            for ws in wb.worksheets:
                output.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    output.append(" | ".join("" if v is None else str(v) for v in row))
            return "\n".join(output)
        if suffix == ".pptx":
            from pptx import Presentation
            deck = Presentation(io.BytesIO(data))
            output = []
            for i, slide in enumerate(deck.slides, start=1):
                output.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        output.append(shape.text)
            return "\n".join(output)
        raise UnsupportedDocumentError(filename)
